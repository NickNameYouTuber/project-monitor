from pptx import Presentation

prs = Presentation('pz.pptx')
slides = prs.slides

print(f'Total slides: {len(slides)}')

# Just count shapes
for i in [21, 22, 23]:
    if i < len(slides):
        s = slides[i]
        print(f'Slide {i+1}: {len(s.shapes)} shapes, layout={s.slide_layout.name}')
