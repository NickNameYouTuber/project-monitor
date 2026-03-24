"""
PowerPoint MCP Server — full-featured presentation editing via python-pptx.
"""

import io
import os
import base64
import json
from pathlib import Path
from typing import Any

from mcp.server.fastmcp import FastMCP
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE

mcp = FastMCP("powerpoint")

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _resolve_path(path: str) -> Path:
    """Resolve and validate a file path."""
    p = Path(path).resolve()
    return p


def _open_pptx(path: str) -> Presentation:
    p = _resolve_path(path)
    if not p.exists():
        raise FileNotFoundError(f"File not found: {p}")
    return Presentation(str(p))


def _save_pptx(prs: Presentation, path: str) -> str:
    p = _resolve_path(path)
    p.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(p))
    return str(p)


def _parse_color(color_str: str | None) -> RGBColor | None:
    if not color_str:
        return None
    color_str = color_str.lstrip("#")
    if len(color_str) == 6:
        return RGBColor(int(color_str[0:2], 16), int(color_str[2:4], 16), int(color_str[4:6], 16))
    return None


def _alignment(align: str | None):
    if not align:
        return None
    mapping = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    return mapping.get(align.lower())


def _get_slide(prs: Presentation, slide_index: int):
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(f"Slide index {slide_index} out of range (0-{len(prs.slides)-1})")
    return prs.slides[slide_index]


def _shape_info(shape) -> dict:
    info: dict[str, Any] = {
        "shape_id": shape.shape_id,
        "name": shape.name,
        "shape_type": str(shape.shape_type),
        "left": shape.left,
        "top": shape.top,
        "width": shape.width,
        "height": shape.height,
    }
    if shape.has_text_frame:
        info["text"] = shape.text_frame.text
        info["paragraphs"] = []
        for para in shape.text_frame.paragraphs:
            p_info = {"text": para.text, "runs": []}
            for run in para.runs:
                r_info = {"text": run.text}
                if run.font.size:
                    r_info["font_size_pt"] = run.font.size.pt
                if run.font.bold is not None:
                    r_info["bold"] = run.font.bold
                if run.font.italic is not None:
                    r_info["italic"] = run.font.italic
                if run.font.color and run.font.color.rgb:
                    r_info["color"] = str(run.font.color.rgb)
                if run.font.name:
                    r_info["font_name"] = run.font.name
                p_info["runs"].append(r_info)
            info["paragraphs"].append(p_info)
    if shape.has_table:
        info["table"] = {
            "rows": shape.table.rows.__len__(),
            "cols": shape.table.columns.__len__(),
        }
    return info


# ---------------------------------------------------------------------------
# Tools — Document management
# ---------------------------------------------------------------------------

@mcp.tool()
def create_presentation(path: str, title: str = "", subtitle: str = "") -> str:
    """Create a new PowerPoint presentation. Optionally adds a title slide.

    Args:
        path: File path for the new .pptx file
        title: Optional title for the first slide
        subtitle: Optional subtitle for the first slide
    """
    prs = Presentation()
    if title:
        layout = prs.slide_layouts[0]  # Title Slide layout
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        if subtitle and slide.placeholders[1]:
            slide.placeholders[1].text = subtitle
    saved = _save_pptx(prs, path)
    return json.dumps({"status": "created", "path": saved, "slides": len(prs.slides)})


@mcp.tool()
def get_presentation_info(path: str) -> str:
    """Get metadata and overview of a PowerPoint presentation.

    Args:
        path: Path to the .pptx file
    """
    prs = _open_pptx(path)
    slides_info = []
    for idx, slide in enumerate(prs.slides):
        s = {
            "index": idx,
            "layout": slide.slide_layout.name,
            "shapes_count": len(slide.shapes),
        }
        if slide.shapes.title:
            s["title"] = slide.shapes.title.text
        slides_info.append(s)

    info = {
        "path": str(_resolve_path(path)),
        "slide_count": len(prs.slides),
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "layouts_available": [l.name for l in prs.slide_layouts],
        "slides": slides_info,
    }
    return json.dumps(info, ensure_ascii=False)


@mcp.tool()
def list_available_presentations(directory: str = ".") -> str:
    """List all .pptx files in a directory.

    Args:
        directory: Directory path to search (default: current directory)
    """
    d = _resolve_path(directory)
    files = list(d.glob("**/*.pptx"))
    return json.dumps([str(f) for f in files], ensure_ascii=False)


@mcp.tool()
def copy_presentation(source_path: str, dest_path: str) -> str:
    """Copy a presentation to a new file.

    Args:
        source_path: Path to the source .pptx file
        dest_path: Path for the copy
    """
    prs = _open_pptx(source_path)
    saved = _save_pptx(prs, dest_path)
    return json.dumps({"status": "copied", "source": str(_resolve_path(source_path)), "dest": saved})


# ---------------------------------------------------------------------------
# Tools — Slide management
# ---------------------------------------------------------------------------

@mcp.tool()
def add_slide(path: str, layout_index: int = 1) -> str:
    """Add a new slide to the presentation.

    Args:
        path: Path to the .pptx file
        layout_index: Index of the slide layout to use (0=Title Slide, 1=Title and Content, 5=Blank, etc.)
    """
    prs = _open_pptx(path)
    if layout_index < 0 or layout_index >= len(prs.slide_layouts):
        return json.dumps({"error": f"Layout index {layout_index} out of range (0-{len(prs.slide_layouts)-1})"})
    layout = prs.slide_layouts[layout_index]
    prs.slides.add_slide(layout)
    _save_pptx(prs, path)
    return json.dumps({
        "status": "added",
        "slide_index": len(prs.slides) - 1,
        "layout": layout.name,
        "total_slides": len(prs.slides),
    })


@mcp.tool()
def delete_slide(path: str, slide_index: int) -> str:
    """Delete a slide from the presentation.

    Args:
        path: Path to the .pptx file
        slide_index: Zero-based index of the slide to delete
    """
    prs = _open_pptx(path)
    if slide_index < 0 or slide_index >= len(prs.slides):
        return json.dumps({"error": f"Slide index {slide_index} out of range"})

    rId = prs.slides._sldIdLst[slide_index].get('r:id')
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[slide_index]
    _save_pptx(prs, path)
    return json.dumps({"status": "deleted", "deleted_index": slide_index, "total_slides": len(prs.slides)})


@mcp.tool()
def duplicate_slide(path: str, slide_index: int) -> str:
    """Duplicate a slide in the presentation (creates a copy right after the original).

    Args:
        path: Path to the .pptx file
        slide_index: Zero-based index of the slide to duplicate
    """
    import copy
    from lxml import etree

    prs = _open_pptx(path)
    slide = _get_slide(prs, slide_index)

    # Add a blank slide with the same layout, then copy XML
    new_slide = prs.slides.add_slide(slide.slide_layout)
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    for shape in slide.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)

    _save_pptx(prs, path)
    return json.dumps({
        "status": "duplicated",
        "source_index": slide_index,
        "new_index": len(prs.slides) - 1,
        "total_slides": len(prs.slides),
    })


@mcp.tool()
def get_slide_info(path: str, slide_index: int) -> str:
    """Get detailed information about a specific slide including all shapes and text.

    Args:
        path: Path to the .pptx file
        slide_index: Zero-based index of the slide
    """
    prs = _open_pptx(path)
    slide = _get_slide(prs, slide_index)

    shapes = [_shape_info(s) for s in slide.shapes]
    info = {
        "slide_index": slide_index,
        "layout": slide.slide_layout.name,
        "shapes": shapes,
    }
    return json.dumps(info, ensure_ascii=False)


# ---------------------------------------------------------------------------
# Tools — Text editing
# ---------------------------------------------------------------------------

@mcp.tool()
def set_text_in_shape(
    path: str,
    slide_index: int,
    shape_name: str,
    text: str,
    font_size: float | None = None,
    font_name: str | None = None,
    bold: bool | None = None,
    italic: bool | None = None,
    color: str | None = None,
    alignment: str | None = None,
) -> str:
    """Set or replace all text in a shape. Clears existing text and sets new text with optional formatting.

    Args:
        path: Path to the .pptx file
        slide_index: Zero-based index of the slide
        shape_name: Name of the shape to edit
        text: New text content
        font_size: Font size in points
        font_name: Font family name
        bold: Whether the text should be bold
        italic: Whether the text should be italic
        color: Hex color string (e.g. "#FF0000")
        alignment: Text alignment (left, center, right, justify)
    """
    prs = _open_pptx(path)
    slide = _get_slide(prs, slide_index)

    shape = None
    for s in slide.shapes:
        if s.name == shape_name:
            shape = s
            break
    if not shape:
        return json.dumps({"error": f"Shape '{shape_name}' not found on slide {slide_index}"})
    if not shape.has_text_frame:
        return json.dumps({"error": f"Shape '{shape_name}' does not have a text frame"})

    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text

    if font_size:
        run.font.size = Pt(font_size)
    if font_name:
        run.font.name = font_name
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    rgb = _parse_color(color)
    if rgb:
        run.font.color.rgb = rgb
    a = _alignment(alignment)
    if a is not None:
        p.alignment = a

    _save_pptx(prs, path)
    return json.dumps({"status": "text_set", "shape": shape_name, "text": text})


@mcp.tool()
def add_paragraph_to_shape(
    path: str,
    slide_index: int,
    shape_name: str,
    text: str,
    font_size: float | None = None,
    font_name: str | None = None,
    bold: bool | None = None,
    italic: bool | None = None,
    color: str | None = None,
    alignment: str | None = None,
    level: int = 0,
) -> str:
    """Add a new paragraph to an existing text frame in a shape.

    Args:
        path: Path to the .pptx file
        slide_index: Zero-based index of the slide
        shape_name: Name of the shape
        text: Text for the new paragraph
        font_size: Font size in points
        font_name: Font family name
        bold: Whether the text should be bold
        italic: Whether the text should be italic
        color: Hex color string (e.g. "#FF0000")
        alignment: Text alignment (left, center, right, justify)
        level: Indentation level (0-8, for bullet lists)
    """
    prs = _open_pptx(path)
    slide = _get_slide(prs, slide_index)

    shape = None
    for s in slide.shapes:
        if s.name == shape_name:
            shape = s
            break
    if not shape or not shape.has_text_frame:
        return json.dumps({"error": f"Shape '{shape_name}' not found or has no text frame"})

    tf = shape.text_frame
    p = tf.add_paragraph()
    p.level = level
    run = p.add_run()
    run.text = text

    if font_size:
        run.font.size = Pt(font_size)
    if font_name:
        run.font.name = font_name
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    rgb = _parse_color(color)
    if rgb:
        run.font.color.rgb = rgb
    a = _alignment(alignment)
    if a is not None:
        p.alignment = a

    _save_pptx(prs, path)
    return json.dumps({"status": "paragraph_added", "shape": shape_name, "text": text})


@mcp.tool()
def find_and_replace_text(path: str, old_text: str, new_text: str, slide_index: int | None = None) -> str:
    """Find and replace text across all slides or a specific slide.

    Args:
        path: Path to the .pptx file
        old_text: Text to find
        new_text: Replacement text
        slide_index: Optional specific slide index (replaces on all slides if omitted)
    """
    prs = _open_pptx(path)
    replacements = 0

    slides = [_get_slide(prs, slide_index)] if slide_index is not None else prs.slides

    for slide in slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if old_text in run.text:
                            run.text = run.text.replace(old_text, new_text)
                            replacements += 1
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if old_text in run.text:
                                    run.text = run.text.replace(old_text, new_text)
                                    replacements += 1

    _save_pptx(prs, path)
    return json.dumps({"status": "replaced", "replacements": replacements, "old_text": old_text, "new_text": new_text})


@mcp.tool()
def get_all_text(path: str, slide_index: int | None = None) -> str:
    """Extract all text from the presentation or a specific slide.

    Args:
        path: Path to the .pptx file
        slide_index: Optional specific slide index
    """
    prs = _open_pptx(path)
    slides = [_get_slide(prs, slide_index)] if slide_index is not None else prs.slides
    result = []
    for idx, slide in enumerate(slides):
        real_idx = slide_index if slide_index is not None else idx
        slide_text = {"slide_index": real_idx, "texts": []}
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text["texts"].append({
                    "shape_name": shape.name,
                    "text": shape.text_frame.text,
                })
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    table_data.append([cell.text for cell in row.cells])
                slide_text["texts"].append({
                    "shape_name": shape.name,
                    "table": table_data,
                })
        result.append(slide_text)
    return json.dumps(result, ensure_ascii=False)


# ---------------------------------------------------------------------------
# Tools — Shapes
# ---------------------------------------------------------------------------

@mcp.tool()
def add_textbox(
    path: str,
    slide_index: int,
    text: str,
    left: float = 1.0,
    top: float = 1.0,
    width: float = 5.0,
    height: float = 1.0,
    font_size: float = 18,
    font_name: str | None = None,
    bold: bool = False,
    italic: bool = False,
    color: str | None = None,
    alignment: str | None = None,
) -> str:
    """Add a text box to a slide.

    Args:
        path: Path to the .pptx file
        slide_index: Zero-based slide index
        text: Text content
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        font_size: Font size in points
        font_name: Font family
        bold: Bold text
        italic: Italic text
        color: Hex color
        alignment: Text alignment
    """
    prs = _open_pptx(path)
    slide = _get_slide(prs, slide_index)

    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)

    if font_name:
        run.font.name = font_name
    if bold:
        run.font.bold = True
    if italic:
        run.font.italic = True
    rgb = _parse_color(color)
    if rgb:
        run.font.color.rgb = rgb
    a = _alignment(alignment)
    if a is not None:
        p.alignment = a

    _save_pptx(prs, path)
    return json.dumps({"status": "textbox_added", "shape_name": txBox.name, "slide_index": slide_index})


@mcp.tool()
def add_shape(
    path: str,
    slide_index: int,
    shape_type: str = "RECTANGLE",
    left: float = 1.0,
    top: float = 1.0,
    width: float = 3.0,
    height: float = 2.0,
    fill_color: str | None = None,
    line_color: str | None = None,
    text: str | None = None,
) -> str:
    """Add an auto-shape (rectangle, oval, etc.) to a slide.

    Args:
        path: Path to the .pptx file
        slide_index: Zero-based slide index
        shape_type: Shape type name (RECTANGLE, OVAL, ROUNDED_RECTANGLE, TRIANGLE, DIAMOND, STAR_5_POINT, ARROW_RIGHT, etc.)
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        fill_color: Fill color as hex string
        line_color: Line/border color as hex string
        text: Optional text inside the shape
    """
    prs = _open_pptx(path)
    slide = _get_slide(prs, slide_index)

    shape_enum = getattr(MSO_SHAPE, shape_type.upper(), None)
    if shape_enum is None:
        available = [s for s in dir(MSO_SHAPE) if not s.startswith("_")]
        return json.dumps({"error": f"Unknown shape type '{shape_type}'", "available": available[:30]})

    shape = slide.shapes.add_shape(
        shape_enum, Inches(left), Inches(top), Inches(width), Inches(height)
    )

    fill_rgb = _parse_color(fill_color)
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb

    line_rgb = _parse_color(line_color)
    if line_rgb:
        shape.line.color.rgb = line_rgb

    if text:
        shape.text = text

    _save_pptx(prs, path)
    return json.dumps({"status": "shape_added", "shape_name": shape.name, "shape_type": shape_type})


@mcp.tool()
def add_image(
    path: str,
    slide_index: int,
    image_path: str,
    left: float = 1.0,
    top: float = 1.0,
    width: float | None = None,
    height: float | None = None,
) -> str:
    """Add an image to a slide.

    Args:
        path: Path to the .pptx file
        slide_index: Zero-based slide index
        image_path: Path to the image file
        left: Left position in inches
        top: Top position in inches
        width: Optional width in inches (maintains aspect ratio if only width or height given)
        height: Optional height in inches
    """
    prs = _open_pptx(path)
    slide = _get_slide(prs, slide_index)

    img_path = _resolve_path(image_path)
    if not img_path.exists():
        return json.dumps({"error": f"Image not found: {img_path}"})

    kwargs: dict[str, Any] = {"left": Inches(left), "top": Inches(top)}
    if width:
        kwargs["width"] = Inches(width)
    if height:
        kwargs["height"] = Inches(height)

    pic = slide.shapes.add_picture(str(img_path), **kwargs)
    _save_pptx(prs, path)
    return json.dumps({"status": "image_added", "shape_name": pic.name, "slide_index": slide_index})


@mcp.tool()
def delete_shape(path: str, slide_index: int, shape_name: str) -> str:
    """Delete a shape from a slide by its name.

    Args:
        path: Path to the .pptx file
        slide_index: Zero-based slide index
        shape_name: Name of the shape to delete
    """
    prs = _open_pptx(path)
    slide = _get_slide(prs, slide_index)

    for shape in slide.shapes:
        if shape.name == shape_name:
            sp = shape._element
            sp.getparent().remove(sp)
            _save_pptx(prs, path)
            return json.dumps({"status": "deleted", "shape_name": shape_name})

    return json.dumps({"error": f"Shape '{shape_name}' not found on slide {slide_index}"})


@mcp.tool()
def move_and_resize_shape(
    path: str,
    slide_index: int,
    shape_name: str,
    left: float | None = None,
    top: float | None = None,
    width: float | None = None,
    height: float | None = None,
) -> str:
    """Move and/or resize a shape on a slide.

    Args:
        path: Path to the .pptx file
        slide_index: Zero-based slide index
        shape_name: Name of the shape
        left: New left position in inches
        top: New top position in inches
        width: New width in inches
        height: New height in inches
    """
    prs = _open_pptx(path)
    slide = _get_slide(prs, slide_index)

    shape = None
    for s in slide.shapes:
        if s.name == shape_name:
            shape = s
            break
    if not shape:
        return json.dumps({"error": f"Shape '{shape_name}' not found"})

    if left is not None:
        shape.left = Inches(left)
    if top is not None:
        shape.top = Inches(top)
    if width is not None:
        shape.width = Inches(width)
    if height is not None:
        shape.height = Inches(height)

    _save_pptx(prs, path)
    return json.dumps({"status": "updated", "shape_name": shape_name})


# ---------------------------------------------------------------------------
# Tools — Tables
# ---------------------------------------------------------------------------

@mcp.tool()
def add_table(
    path: str,
    slide_index: int,
    rows: int,
    cols: int,
    data: list[list[str]] | None = None,
    left: float = 1.0,
    top: float = 2.0,
    width: float = 8.0,
    height: float = 3.0,
) -> str:
    """Add a table to a slide.

    Args:
        path: Path to the .pptx file
        slide_index: Zero-based slide index
        rows: Number of rows
        cols: Number of columns
        data: Optional 2D list of cell text values
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
    """
    prs = _open_pptx(path)
    slide = _get_slide(prs, slide_index)

    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table

    if data:
        for r_idx, row_data in enumerate(data):
            if r_idx >= rows:
                break
            for c_idx, cell_text in enumerate(row_data):
                if c_idx >= cols:
                    break
                table.cell(r_idx, c_idx).text = str(cell_text)

    _save_pptx(prs, path)
    return json.dumps({"status": "table_added", "shape_name": table_shape.name, "rows": rows, "cols": cols})


@mcp.tool()
def set_table_cell(
    path: str,
    slide_index: int,
    shape_name: str,
    row: int,
    col: int,
    text: str,
    font_size: float | None = None,
    bold: bool | None = None,
    color: str | None = None,
) -> str:
    """Set text and formatting for a specific table cell.

    Args:
        path: Path to the .pptx file
        slide_index: Zero-based slide index
        shape_name: Name of the table shape
        row: Zero-based row index
        col: Zero-based column index
        text: Cell text
        font_size: Font size in points
        bold: Bold text
        color: Hex color
    """
    prs = _open_pptx(path)
    slide = _get_slide(prs, slide_index)

    shape = None
    for s in slide.shapes:
        if s.name == shape_name and s.has_table:
            shape = s
            break
    if not shape:
        return json.dumps({"error": f"Table '{shape_name}' not found on slide {slide_index}"})

    cell = shape.table.cell(row, col)
    cell.text = text

    if font_size or bold is not None or color:
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                if font_size:
                    run.font.size = Pt(font_size)
                if bold is not None:
                    run.font.bold = bold
                rgb = _parse_color(color)
                if rgb:
                    run.font.color.rgb = rgb

    _save_pptx(prs, path)
    return json.dumps({"status": "cell_set", "row": row, "col": col, "text": text})


# ---------------------------------------------------------------------------
# Tools — Design & Formatting
# ---------------------------------------------------------------------------

@mcp.tool()
def set_slide_background_color(path: str, slide_index: int, color: str) -> str:
    """Set the background color of a slide.

    Args:
        path: Path to the .pptx file
        slide_index: Zero-based slide index
        color: Hex color string (e.g. "#1A1A2E")
    """
    prs = _open_pptx(path)
    slide = _get_slide(prs, slide_index)

    background = slide.background
    fill = background.fill
    fill.solid()
    rgb = _parse_color(color)
    if rgb:
        fill.fore_color.rgb = rgb

    _save_pptx(prs, path)
    return json.dumps({"status": "background_set", "slide_index": slide_index, "color": color})


@mcp.tool()
def set_slide_background_image(path: str, slide_index: int, image_path: str) -> str:
    """Set a background image for a slide.

    Args:
        path: Path to the .pptx file
        slide_index: Zero-based slide index
        image_path: Path to the image file
    """
    from pptx.oxml.ns import qn
    import copy

    prs = _open_pptx(path)
    slide = _get_slide(prs, slide_index)
    img_path = _resolve_path(image_path)

    if not img_path.exists():
        return json.dumps({"error": f"Image not found: {img_path}"})

    # Add image as background via relationship
    rel = slide.part.get_or_add_image_part(str(img_path))
    rId = rel[0]  # relationship id

    bg = slide.background
    bg_elem = bg._element
    bgPr = bg_elem.find(qn('p:bgPr'))
    if bgPr is None:
        from lxml import etree
        bgPr = etree.SubElement(bg_elem, qn('p:bgPr'))

    # Clear existing fill
    for child in list(bgPr):
        bgPr.remove(child)

    # Add blipFill
    from lxml import etree
    blipFill = etree.SubElement(bgPr, qn('a:blipFill'))
    blip = etree.SubElement(blipFill, qn('a:blip'))
    blip.set(qn('r:embed'), rId)
    stretch = etree.SubElement(blipFill, qn('a:stretch'))
    etree.SubElement(stretch, qn('a:fillRect'))
    etree.SubElement(bgPr, qn('a:effectLst'))

    _save_pptx(prs, path)
    return json.dumps({"status": "background_image_set", "slide_index": slide_index})


@mcp.tool()
def format_shape_fill(
    path: str,
    slide_index: int,
    shape_name: str,
    fill_color: str | None = None,
    transparency: float = 0.0,
    no_fill: bool = False,
) -> str:
    """Set fill properties of a shape.

    Args:
        path: Path to the .pptx file
        slide_index: Zero-based slide index
        shape_name: Shape name
        fill_color: Hex color for solid fill
        transparency: Fill transparency (0.0 to 1.0)
        no_fill: If True, remove fill completely
    """
    prs = _open_pptx(path)
    slide = _get_slide(prs, slide_index)

    shape = None
    for s in slide.shapes:
        if s.name == shape_name:
            shape = s
            break
    if not shape:
        return json.dumps({"error": f"Shape '{shape_name}' not found"})

    if no_fill:
        shape.fill.background()
    elif fill_color:
        shape.fill.solid()
        rgb = _parse_color(fill_color)
        if rgb:
            shape.fill.fore_color.rgb = rgb

    _save_pptx(prs, path)
    return json.dumps({"status": "fill_updated", "shape_name": shape_name})


@mcp.tool()
def set_shape_border(
    path: str,
    slide_index: int,
    shape_name: str,
    color: str | None = None,
    width: float | None = None,
) -> str:
    """Set border/line properties of a shape.

    Args:
        path: Path to the .pptx file
        slide_index: Zero-based slide index
        shape_name: Shape name
        color: Border color as hex string
        width: Border width in points
    """
    prs = _open_pptx(path)
    slide = _get_slide(prs, slide_index)

    shape = None
    for s in slide.shapes:
        if s.name == shape_name:
            shape = s
            break
    if not shape:
        return json.dumps({"error": f"Shape '{shape_name}' not found"})

    if color:
        rgb = _parse_color(color)
        if rgb:
            shape.line.color.rgb = rgb
    if width is not None:
        shape.line.width = Pt(width)

    _save_pptx(prs, path)
    return json.dumps({"status": "border_updated", "shape_name": shape_name})


@mcp.tool()
def set_slide_dimensions(path: str, width_inches: float, height_inches: float) -> str:
    """Set the slide dimensions for the entire presentation.

    Args:
        path: Path to the .pptx file
        width_inches: Slide width in inches (standard 16:9 = 13.333, 4:3 = 10.0)
        height_inches: Slide height in inches (standard 16:9 = 7.5, 4:3 = 7.5)
    """
    prs = _open_pptx(path)
    prs.slide_width = Inches(width_inches)
    prs.slide_height = Inches(height_inches)
    _save_pptx(prs, path)
    return json.dumps({"status": "dimensions_set", "width": width_inches, "height": height_inches})


@mcp.tool()
def apply_theme_colors(
    path: str,
    slide_index: int,
    background_color: str | None = None,
    title_color: str | None = None,
    body_color: str | None = None,
    accent_color: str | None = None,
) -> str:
    """Apply a color theme to a slide — sets background and text colors for a cohesive look.

    Args:
        path: Path to the .pptx file
        slide_index: Zero-based slide index
        background_color: Hex color for slide background
        title_color: Hex color for title text
        body_color: Hex color for body/content text
        accent_color: Hex color for accent shapes
    """
    prs = _open_pptx(path)
    slide = _get_slide(prs, slide_index)

    if background_color:
        fill = slide.background.fill
        fill.solid()
        rgb = _parse_color(background_color)
        if rgb:
            fill.fore_color.rgb = rgb

    for shape in slide.shapes:
        if shape.has_text_frame:
            is_title = shape == slide.shapes.title or "title" in shape.name.lower()
            color_str = title_color if is_title else body_color
            if color_str:
                rgb = _parse_color(color_str)
                if rgb:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.color.rgb = rgb

    _save_pptx(prs, path)
    return json.dumps({"status": "theme_applied", "slide_index": slide_index})


@mcp.tool()
def apply_template(path: str, template_path: str) -> str:
    """Create a new presentation based on a template .pptx file (copies the template).

    Args:
        path: Path for the new presentation
        template_path: Path to the template .pptx file
    """
    template = _resolve_path(template_path)
    if not template.exists():
        return json.dumps({"error": f"Template not found: {template}"})

    prs = Presentation(str(template))
    saved = _save_pptx(prs, path)
    return json.dumps({"status": "created_from_template", "path": saved, "template": str(template)})


# ---------------------------------------------------------------------------
# Tools — Notes
# ---------------------------------------------------------------------------

@mcp.tool()
def set_slide_notes(path: str, slide_index: int, notes: str) -> str:
    """Set speaker notes for a slide.

    Args:
        path: Path to the .pptx file
        slide_index: Zero-based slide index
        notes: Speaker notes text
    """
    prs = _open_pptx(path)
    slide = _get_slide(prs, slide_index)

    if not slide.has_notes_slide:
        slide.notes_slide  # creates notes slide
    slide.notes_slide.notes_text_frame.text = notes

    _save_pptx(prs, path)
    return json.dumps({"status": "notes_set", "slide_index": slide_index})


@mcp.tool()
def get_slide_notes(path: str, slide_index: int) -> str:
    """Get speaker notes from a slide.

    Args:
        path: Path to the .pptx file
        slide_index: Zero-based slide index
    """
    prs = _open_pptx(path)
    slide = _get_slide(prs, slide_index)

    notes = ""
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text

    return json.dumps({"slide_index": slide_index, "notes": notes})


# ---------------------------------------------------------------------------
# Tools — Placeholder management
# ---------------------------------------------------------------------------

@mcp.tool()
def list_placeholders(path: str, slide_index: int) -> str:
    """List all placeholders on a slide with their indices and types.

    Args:
        path: Path to the .pptx file
        slide_index: Zero-based slide index
    """
    prs = _open_pptx(path)
    slide = _get_slide(prs, slide_index)

    placeholders = []
    for ph in slide.placeholders:
        info = {
            "idx": ph.placeholder_format.idx,
            "name": ph.name,
            "type": str(ph.placeholder_format.type),
            "width": ph.width,
            "height": ph.height,
        }
        if ph.has_text_frame:
            info["text"] = ph.text_frame.text
        placeholders.append(info)

    return json.dumps(placeholders, ensure_ascii=False)


@mcp.tool()
def set_placeholder_text(
    path: str,
    slide_index: int,
    placeholder_idx: int,
    text: str,
    font_size: float | None = None,
    font_name: str | None = None,
    bold: bool | None = None,
    color: str | None = None,
) -> str:
    """Set text in a slide placeholder by its index.

    Args:
        path: Path to the .pptx file
        slide_index: Zero-based slide index
        placeholder_idx: Placeholder index
        text: New text
        font_size: Font size in points
        font_name: Font family
        bold: Bold text
        color: Hex color
    """
    prs = _open_pptx(path)
    slide = _get_slide(prs, slide_index)

    if placeholder_idx not in [ph.placeholder_format.idx for ph in slide.placeholders]:
        return json.dumps({"error": f"Placeholder {placeholder_idx} not found"})

    ph = slide.placeholders[placeholder_idx]
    ph.text = text

    if font_size or font_name or bold is not None or color:
        for para in ph.text_frame.paragraphs:
            for run in para.runs:
                if font_size:
                    run.font.size = Pt(font_size)
                if font_name:
                    run.font.name = font_name
                if bold is not None:
                    run.font.bold = bold
                rgb = _parse_color(color)
                if rgb:
                    run.font.color.rgb = rgb

    _save_pptx(prs, path)
    return json.dumps({"status": "placeholder_set", "idx": placeholder_idx, "text": text})


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def main():
    mcp.run(transport="stdio")


if __name__ == "__main__":
    main()
